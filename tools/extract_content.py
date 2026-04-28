#!/usr/bin/env python3
"""
extract_content.py — Source document → structured training content JSON.

Takes a PDF, DOCX, or PPTX source file (training material, policy doc,
publication, slide deck) and outputs a structured JSON file that the
training build script can consume.

Usage:
    python3 extract_content.py <source_file> [-o output.json]

Example:
    python3 extract_content.py irs_pub_1075.pdf -o content.json
    python3 extract_content.py existing_training.pptx -o content.json
    python3 extract_content.py policy.docx -o content.json

Output JSON structure:
    {
      "source": { "filename": "...", "type": "pdf|docx|pptx", "pages": N },
      "title": "Inferred document title",
      "sections": [
        {
          "level": 1,           # 1 = top section, 2 = subsection, etc.
          "heading": "Section title",
          "body": ["Paragraph 1", "Paragraph 2", ...],
          "bullets": ["Bullet 1", "Bullet 2", ...],
          "page": N,            # source page or slide number
          "notes": "Speaker notes if PPTX"
        },
        ...
      ]
    }

The output is *structural* — what's in the source document, organized.
It is NOT yet training content (no scenarios, no rewriting). The next
step is human review + edits to tailor it to the training format.
"""

import argparse
import json
import os
import re
import sys
from pathlib import Path


# ============================================================================
# PDF extraction
# ============================================================================

def extract_pdf(path):
    """Pull text + structure from a PDF using pypdf, with heading detection."""
    from pypdf import PdfReader

    reader = PdfReader(path)
    page_count = len(reader.pages)

    sections = []
    current_section = None

    for page_num, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception as e:
            print(f"  warning: failed to extract page {page_num}: {e}", file=sys.stderr)
            continue

        # Split by lines and look for heading patterns
        lines = [ln.rstrip() for ln in text.split("\n")]
        for line in lines:
            stripped = line.strip()
            if not stripped:
                continue

            # Heading detection: numbered sections like "1.0", "2.5", "3.4.1"
            # or ALL CAPS lines, or short lines that look like headings
            heading_level = detect_heading_level(stripped)

            if heading_level:
                if current_section:
                    sections.append(current_section)
                current_section = {
                    "level": heading_level,
                    "heading": stripped,
                    "body": [],
                    "bullets": [],
                    "page": page_num,
                }
            else:
                if current_section is None:
                    current_section = {
                        "level": 0,
                        "heading": "(Document Preamble)",
                        "body": [],
                        "bullets": [],
                        "page": page_num,
                    }
                # Bullet detection
                if re.match(r"^[\*\-\u2022\u00b7]\s+", stripped) or re.match(r"^\(?[a-z0-9]\)\s+", stripped):
                    bullet_text = re.sub(r"^[\*\-\u2022\u00b7]\s+", "", stripped)
                    bullet_text = re.sub(r"^\(?[a-z0-9]\)\s+", "", bullet_text)
                    current_section["bullets"].append(bullet_text)
                else:
                    current_section["body"].append(stripped)

    if current_section:
        sections.append(current_section)

    # Join body lines that look like wrapped paragraphs back together
    for section in sections:
        section["body"] = merge_wrapped_lines(section["body"])

    title = infer_pdf_title(reader, sections)

    return {
        "source": {
            "filename": os.path.basename(path),
            "type": "pdf",
            "pages": page_count,
        },
        "title": title,
        "sections": sections,
    }


def detect_heading_level(line):
    """Return heading level 1-4, or 0 if not a heading."""
    if len(line) > 200:
        return 0  # too long to be a heading

    # Numbered section patterns
    m = re.match(r"^(\d+)\.0\s+[A-Z]", line)
    if m:
        return 1
    m = re.match(r"^(\d+)\.(\d+)\s+[A-Z]", line)
    if m:
        return 2
    m = re.match(r"^(\d+)\.(\d+)\.(\d+)\s+[A-Z]", line)
    if m:
        return 3

    # ALL CAPS heading (short)
    if len(line) < 80 and line == line.upper() and re.search(r"[A-Z]", line) and len(line.split()) >= 2:
        return 1

    return 0


def merge_wrapped_lines(lines):
    """Merge lines that look like wrapped paragraph fragments."""
    merged = []
    buffer = ""
    for line in lines:
        if not line:
            if buffer:
                merged.append(buffer.strip())
                buffer = ""
            continue
        # If buffer is non-empty and line continues a sentence, append
        if buffer and not buffer.rstrip().endswith((".", "!", "?", ":", ";")):
            buffer += " " + line
        else:
            if buffer:
                merged.append(buffer.strip())
            buffer = line
    if buffer:
        merged.append(buffer.strip())
    return merged


def infer_pdf_title(reader, sections):
    """Pull a sensible title from the PDF metadata or first heading."""
    if reader.metadata and reader.metadata.title:
        return reader.metadata.title.strip()
    for section in sections:
        if section["level"] == 1:
            return section["heading"]
    return "Untitled Document"


# ============================================================================
# DOCX extraction
# ============================================================================

def extract_docx(path):
    """Pull structured content from a DOCX using python-docx."""
    from docx import Document

    doc = Document(path)
    sections = []
    current_section = None

    title = doc.core_properties.title or "Untitled Document"

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        style_name = para.style.name if para.style else ""

        # Heading detection from style
        heading_level = 0
        if style_name.startswith("Heading "):
            try:
                heading_level = int(style_name.replace("Heading ", ""))
            except ValueError:
                heading_level = 1
        elif style_name == "Title":
            heading_level = 1
            title = text

        if heading_level:
            if current_section:
                sections.append(current_section)
            current_section = {
                "level": heading_level,
                "heading": text,
                "body": [],
                "bullets": [],
                "page": None,  # DOCX doesn't expose page numbers reliably
            }
        else:
            if current_section is None:
                current_section = {
                    "level": 0,
                    "heading": "(Document Preamble)",
                    "body": [],
                    "bullets": [],
                    "page": None,
                }

            # Bullet detection from list style
            if "List" in style_name or "Bullet" in style_name:
                current_section["bullets"].append(text)
            else:
                current_section["body"].append(text)

    if current_section:
        sections.append(current_section)

    # Also extract tables — flatten as bullet rows for now
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text and current_section:
                current_section["bullets"].append(row_text)

    return {
        "source": {
            "filename": os.path.basename(path),
            "type": "docx",
            "pages": None,
        },
        "title": title,
        "sections": sections,
    }


# ============================================================================
# PPTX extraction
# ============================================================================

def extract_pptx(path):
    """Pull each slide's content + speaker notes from a PPTX."""
    from pptx import Presentation

    prs = Presentation(path)
    sections = []
    title = "Untitled Presentation"

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_title = ""
        body_parts = []
        bullets = []
        notes = ""

        # Pull title and body from placeholders
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            for para_idx, para in enumerate(tf.paragraphs):
                text = para.text.strip()
                if not text:
                    continue

                # First paragraph of title placeholder = slide title
                if shape.is_placeholder and shape.placeholder_format and shape.placeholder_format.idx == 0:
                    if not slide_title:
                        slide_title = text
                        if slide_idx == 1:
                            title = text
                        continue

                # Bulleted content (level > 0 or list-y formatting)
                if para.level > 0:
                    bullets.append(text)
                else:
                    body_parts.append(text)

        # Pull speaker notes
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            notes = notes_tf.text.strip() if notes_tf else ""

        section = {
            "level": 1,
            "heading": slide_title or f"Slide {slide_idx}",
            "body": body_parts,
            "bullets": bullets,
            "page": slide_idx,  # slide number
        }
        if notes:
            section["notes"] = notes
        sections.append(section)

    return {
        "source": {
            "filename": os.path.basename(path),
            "type": "pptx",
            "pages": len(prs.slides),
        },
        "title": title,
        "sections": sections,
    }


# ============================================================================
# Main dispatch
# ============================================================================

def extract(path):
    """Dispatch to the right extractor based on extension."""
    path = Path(path)
    if not path.exists():
        raise FileNotFoundError(f"Source file not found: {path}")

    ext = path.suffix.lower()
    if ext == ".pdf":
        return extract_pdf(str(path))
    elif ext == ".docx":
        return extract_docx(str(path))
    elif ext == ".pptx":
        return extract_pptx(str(path))
    else:
        raise ValueError(f"Unsupported file type: {ext}. Use .pdf, .docx, or .pptx.")


def main():
    parser = argparse.ArgumentParser(description="Extract structured content from a source document for training adaptation.")
    parser.add_argument("source", help="Path to source PDF, DOCX, or PPTX file")
    parser.add_argument("-o", "--output", default="content.json", help="Output JSON path (default: content.json)")
    parser.add_argument("--summary", action="store_true", help="Print a summary instead of full JSON to stdout")
    args = parser.parse_args()

    print(f"Extracting from {args.source}...", file=sys.stderr)
    result = extract(args.source)

    print(f"  Found {len(result['sections'])} sections", file=sys.stderr)
    print(f"  Title: {result['title']}", file=sys.stderr)

    with open(args.output, "w", encoding="utf-8") as f:
        json.dump(result, f, indent=2, ensure_ascii=False)
    print(f"Wrote {args.output}", file=sys.stderr)

    if args.summary:
        print(f"\nDocument: {result['title']}")
        print(f"Type: {result['source']['type']}, pages/slides: {result['source']['pages']}")
        print(f"Sections: {len(result['sections'])}\n")
        for i, s in enumerate(result["sections"][:20], 1):
            indent = "  " * (s["level"])
            print(f"{indent}[{s.get('page','?')}] {s['heading'][:80]}")
            if s["body"]:
                preview = s["body"][0][:100] if s["body"][0] else ""
                if preview:
                    print(f"{indent}    {preview}...")
        if len(result["sections"]) > 20:
            print(f"\n  ... and {len(result['sections']) - 20} more sections")


if __name__ == "__main__":
    main()
